"""
Vital Guardian — FYP-2 Evaluation Slides Generator
Run: venv\Scripts\python scripts/generate_slides.py
Output: Vital_Guardian_FYP2_Slides.pptx  (in project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF9, 0xFA)
DARK       = RGBColor(0x1A, 0x1D, 0x23)
ACCENT     = RGBColor(0x36, 0x6C, 0xF0)   # blue
MID        = RGBColor(0x55, 0x65, 0x7A)   # secondary text
LIGHT_BG   = RGBColor(0xF1, 0xF4, 0xF9)   # card fill
SUCCESS    = RGBColor(0x16, 0xA3, 0x4A)
WARNING    = RGBColor(0xD9, 0x77, 0x06)
DANGER     = RGBColor(0xDC, 0x26, 0x26)
INFO       = RGBColor(0x25, 0x63, 0xEB)
RULE       = RGBColor(0xD1, 0xD8, 0xE4)   # divider

# ── Layout constants ─────────────────────────────────────────────────────────
W  = Inches(13.33)   # widescreen 16:9 width
H  = Inches(7.50)    # widescreen 16:9 height
ML = Inches(0.70)    # left margin
MT = Inches(0.60)    # top margin for content
CW = Inches(11.93)   # content width


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = OFF_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=Pt(14), bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txbox


def accent_bar(slide, left=ML, top=MT, width=Inches(0.45), height=Pt(4)):
    rect(slide, left, top, width, height, fill_color=ACCENT)


def section_label(slide, text, top):
    add_text(slide, text.upper(), ML, top, CW, Inches(0.30),
             size=Pt(9), bold=True, color=MID)


def slide_title(slide, title, subtitle=None, top=MT):
    add_text(slide, title, ML, top, CW, Inches(0.55),
             size=Pt(30), bold=True, color=DARK)
    if subtitle:
        add_text(slide, subtitle, ML, top + Inches(0.55), CW, Inches(0.35),
                 size=Pt(13), color=MID)


def rule(slide, top):
    rect(slide, ML, top, CW, Pt(1), fill_color=RULE)


def add_table(slide, headers, rows,
              left=ML, top=Inches(2.0), width=CW, row_h=Inches(0.36),
              col_widths=None, header_fill=DARK, row_fills=None):
    """Draw a simple table using rectangles and text boxes."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    header_h = Inches(0.38)
    cur_top = top

    # Header row
    x = left
    for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
        rect(slide, x, cur_top, cw, header_h, fill_color=header_fill)
        add_text(slide, hdr, x + Pt(6), cur_top + Pt(4), cw - Pt(8), header_h,
                 size=Pt(10), bold=True, color=WHITE)
        x += cw
    cur_top += header_h

    # Data rows
    for ri, row in enumerate(rows):
        fill = LIGHT_BG if ri % 2 == 0 else WHITE
        if row_fills and ri < len(row_fills) and row_fills[ri]:
            fill = row_fills[ri]
        x = left
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            rect(slide, x, cur_top, cw, row_h,
                 fill_color=fill, line_color=RULE, line_width=Pt(0.5))
            add_text(slide, str(cell), x + Pt(6), cur_top + Pt(3), cw - Pt(8), row_h,
                     size=Pt(10), color=DARK)
            x += cw
        cur_top += row_h

    return cur_top   # returns bottom y for layout chaining


def stat_card(slide, value, label, left, top, width=Inches(2.60), height=Inches(0.95),
              value_color=ACCENT):
    rect(slide, left, top, width, height, fill_color=LIGHT_BG)
    add_text(slide, value, left + Pt(10), top + Pt(6), width - Pt(12), Inches(0.42),
             size=Pt(22), bold=True, color=value_color)
    add_text(slide, label, left + Pt(10), top + Inches(0.50), width - Pt(12), Inches(0.30),
             size=Pt(10), color=MID)


def info_card(slide, title, bullets, left, top, width=Inches(5.70), height=Inches(1.80),
              title_color=DARK):
    rect(slide, left, top, width, height, fill_color=LIGHT_BG)
    rect(slide, left, top, Pt(4), height, fill_color=ACCENT)
    add_text(slide, title, left + Pt(12), top + Pt(8), width - Pt(16), Inches(0.28),
             size=Pt(11), bold=True, color=title_color)
    bullet_top = top + Inches(0.34)
    bullet_h   = (height - Inches(0.38)) / max(len(bullets), 1)
    for b in bullets:
        add_text(slide, f"  {b}", left + Pt(12), bullet_top, width - Pt(16), bullet_h,
                 size=Pt(10), color=MID)
        bullet_top += bullet_h


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def s1_title(prs):
    sl = blank_slide(prs); bg(sl, WHITE)

    # Left accent bar (full height)
    rect(sl, 0, 0, Inches(0.08), H, fill_color=ACCENT)

    # Top label
    add_text(sl, "FYP-2 EVALUATION  ·  APRIL 2026",
             ML, Inches(0.90), CW, Inches(0.28),
             size=Pt(9), bold=True, color=MID)

    # Title
    add_text(sl, "Vital Guardian", ML, Inches(1.32), CW, Inches(0.80),
             size=Pt(42), bold=True, color=DARK)

    # Subtitle
    add_text(sl, "AI-Based Real-Time Patient Monitoring System",
             ML, Inches(2.18), CW, Inches(0.40),
             size=Pt(17), color=MID)

    rule(sl, Inches(2.75))

    # Team cards
    card_data = [
        ("Mahad Khan",   "Vision Module + UI / Dashboard"),
        ("Eman Hassan",  "LLM / Cognitive Core Module"),
        ("Hasan Waqar",  "Auditory Watchdog Module"),
    ]
    cw = Inches(3.80); gap = Inches(0.12)
    x = ML
    for name, role in card_data:
        rect(sl, x, Inches(3.05), cw, Inches(0.88), fill_color=LIGHT_BG)
        add_text(sl, name, x + Pt(10), Inches(3.12), cw - Pt(14), Inches(0.32),
                 size=Pt(12), bold=True, color=DARK)
        add_text(sl, role, x + Pt(10), Inches(3.44), cw - Pt(14), Inches(0.36),
                 size=Pt(10), color=MID)
        x += cw + gap

    # Supervisor / university
    add_text(sl, "Supervisor:  Dr. Basharat Hussain",
             ML, Inches(4.22), CW, Inches(0.28), size=Pt(12), color=DARK)
    add_text(sl, "FAST NUCES Islamabad  ·  Department of Computer Science",
             ML, Inches(4.55), CW, Inches(0.28), size=Pt(11), color=MID)


def s2_overview(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Project Overview",
                "What we are building and why it matters", top=MT + Pt(10))

    # Problem statement box
    rect(sl, ML, Inches(1.55), CW, Inches(0.92), fill_color=LIGHT_BG)
    rect(sl, ML, Inches(1.55), Pt(4), Inches(0.92), fill_color=DANGER)
    add_text(sl, "Problem Statement", ML + Pt(12), Inches(1.60), CW - Pt(16), Inches(0.28),
             size=Pt(10), bold=True, color=DANGER)
    add_text(sl,
             "Hospital wards lack continuous automated monitoring. Nurses cannot watch every patient "
             "simultaneously — resulting in delayed detection of falls, seizures, and acute respiratory distress. "
             "Every minute of delay directly impacts patient outcomes.",
             ML + Pt(12), Inches(1.90), CW - Pt(16), Inches(0.55),
             size=Pt(10), color=MID)

    section_label(sl, "System Objectives", Inches(2.62))
    add_table(sl,
              ["#", "Objective", "Technology Used"],
              [
                  ["1", "Detect patient falls in real time",             "YOLO11n (OpenVINO) + MoViNet-A2"],
                  ["2", "Detect seizure episodes in real time",          "YOLO11n (OpenVINO) + MoViNet-A2"],
                  ["3", "Detect respiratory distress from audio",        "YAMNet (TF Hub) + Silero VAD"],
                  ["4", "Transcribe patient speech for keyword alerts",  "Faster-Whisper Tiny — English + Urdu"],
                  ["5", "Verify and clinically enrich alerts via LLM",   "Gemini 1.5 Flash — two-tier pipeline"],
                  ["6", "Stream all incidents to a clinical dashboard",  "FastAPI + WebSocket + PostgreSQL"],
              ],
              top=Inches(2.88), col_widths=[Inches(0.40), Inches(6.00), Inches(5.53)])


def s3_feedback(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Faculty Feedback Addressed",
                "FYP-1 panel evaluation and actions taken since")

    # Verdict box
    rect(sl, ML, Inches(1.55), CW, Inches(0.88), fill_color=RGBColor(0xDC, 0xFC, 0xE7))
    rect(sl, ML, Inches(1.55), Pt(4), Inches(0.88), fill_color=SUCCESS)
    add_text(sl, "FYP-1 Panel Verdict — POSITIVE",
             ML + Pt(12), Inches(1.61), CW - Pt(16), Inches(0.28),
             size=Pt(10), bold=True, color=SUCCESS)
    add_text(sl,
             '"Good work — proceed as per plan."  No corrective action items were issued.',
             ML + Pt(12), Inches(1.91), CW - Pt(16), Inches(0.36),
             size=Pt(10), color=MID)

    section_label(sl, "Deliverables Completed Post FYP-1", Inches(2.56))
    add_table(sl,
              ["Area", "What was delivered"],
              [
                  ["System Integration",    "All three modules connected into a single running FastAPI pipeline"],
                  ["Cognitive Core",        "Two-tier Gemini pipeline: binary verification + clinical enrichment"],
                  ["Auditory Watchdog",     "Integrated with demo server — rolling accumulator, one-shot Gemini gate, clip injection"],
                  ["Dashboard",             "Production-grade real-time UI with WebSocket streaming and alert log"],
                  ["Database",              "PostgreSQL: incident log, audit trail, nurse/patient CRUD, RBAC"],
                  ["Model Refinement",      "Threshold tuning, top-2 confidence smoothing, pre-buffer warm-up for MoViNet"],
                  ["Demo Pipeline",         "Multi-patient demo: 7 profiles, synchronized audio playback, live feed mode"],
              ],
              top=Inches(2.82), col_widths=[Inches(2.50), Inches(9.43)])


def s4_progress(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Progress Since FYP-1",
                "Component-by-component comparison")

    add_table(sl,
              ["Component", "At FYP-1", "Now (FYP-2)"],
              [
                  ["Vision Module",       "YOLO + MoViNet trained in isolation",          "Real-time pipeline — pre-buffering, FPS normalisation, top-2 smoothing"],
                  ["Auditory Watchdog",   "YAMNet + Whisper trained, not connected",       "Fully integrated — Privacy Shield, accumulator, one-shot Gemini gate"],
                  ["System Integration",  "All modules separate, no end-to-end flow",      "Single FastAPI process — parallel vision + audio with async WS broadcast"],
                  ["Cognitive Core",      "Architecture planned, not implemented",          "Two-tier Gemini: auto-confirm bypass, binary T2, clinical enrichment T3"],
                  ["Web Dashboard",       "Not started",                                   "Live streaming dashboard — gauges, alert feed, Cognitive Core card"],
                  ["Database",            "Not started",                                   "PostgreSQL — incident log, audit trail, nurse/patient CRUD, RBAC"],
                  ["Demo Pipeline",       "Not started",                                   "7 patient profiles, audio clip sync, live feed mode, post-alert review hold"],
              ],
              top=Inches(1.55), col_widths=[Inches(2.10), Inches(4.20), Inches(5.63)])

    # Stats strip
    stats = [("7", "Patient profiles"), ("3", "AI modules integrated"), ("2-tier", "Gemini pipeline")]
    sw = Inches(3.70); sx = ML
    for val, lbl in stats:
        stat_card(sl, val, lbl, sx, Inches(6.30), width=sw, height=Inches(0.90), value_color=ACCENT)
        sx += sw + Inches(0.12)


def s5_architecture(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "System Architecture",
                "Four modules, one pipeline — vision and audio run in parallel, both feed the Cognitive Core")

    cards = [
        ("Visual Guardian", "Vision", ACCENT, [
            "Person Detection:  YOLO11n — OpenVINO IR, CPU, every frame",
            "Activity Classification:  MoViNet-A2 — 32 frames (fall) · 64 frames (seizure)",
            "Smoothing:  Top-2 confidence mean + SegmentConsolidator streak gate",
            "Events:  Fall · Seizure · Normal activity",
        ]),
        ("Auditory Watchdog", "Audio", WARNING, [
            "Privacy Shield:  Silero VAD (ONNX) — visitor mode, sentence flush",
            "Distress Classifier:  YAMNet — 18 respiratory classes, 3-tier severity",
            "Keyword Spotter:  Faster-Whisper Tiny — English + Urdu, hallucination filtered",
            "Accumulator:  15-second rolling score · one-shot Gemini gate per patient",
        ]),
        ("Cognitive Core", "LLM", SUCCESS, [
            "Model:  Gemini 1.5 Flash (multimodal — text + video frames)",
            "Tier 2 — Binary:  CONFIRMED / SUPPRESSED + reasoning in ~1-2 s",
            "Tier 3 — Clinical:  Narrative · severity rating · nursing action plan",
            "Auto-confirm bypass:  ML confidence >= 50% skips Tier 2 directly",
        ]),
        ("Web Dashboard", "Infra", MID, [
            "Backend:  FastAPI + WebSocket streaming (uvicorn)",
            "Frontend:  Vanilla JS — AudioContext synthesizer, Web Speech API TTS",
            "Database:  PostgreSQL / SQLAlchemy — incident log, audit trail, RBAC",
            "Auth:  Token-based login · Nurse and Admin role separation",
        ]),
    ]

    cw = Inches(5.84); gap = Inches(0.25)
    positions = [
        (ML,           Inches(1.55)),
        (ML + cw + gap, Inches(1.55)),
        (ML,           Inches(4.00)),
        (ML + cw + gap, Inches(4.00)),
    ]
    for (left, top), (title, tag, color, bullets) in zip(positions, cards):
        rect(sl, left, top, cw, Inches(2.22), fill_color=LIGHT_BG)
        rect(sl, left, top, Pt(4), Inches(2.22), fill_color=color)
        add_text(sl, title, left + Pt(12), top + Pt(7), cw - Pt(70), Inches(0.28),
                 size=Pt(11), bold=True, color=DARK)
        # tag pill
        rect(sl, left + cw - Inches(0.85), top + Pt(6), Inches(0.78), Pt(18),
             fill_color=color)
        add_text(sl, tag, left + cw - Inches(0.82), top + Pt(7), Inches(0.74), Pt(16),
                 size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        by = top + Inches(0.36)
        for b in bullets:
            add_text(sl, f"- {b}", left + Pt(12), by, cw - Pt(16), Inches(0.36),
                     size=Pt(9.5), color=MID)
            by += Inches(0.40)


def s6_testing(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Testing, Integration & Quality",
                "Validation strategies at module level and across the full pipeline")

    # Left table
    section_label(sl, "Module-Level Testing", Inches(1.55))
    add_table(sl,
              ["Module", "Validation Method"],
              [
                  ["Visual Guardian",   "Held-out unseen test sets (950 fall · 186 seizure clips)"],
                  ["Auditory Watchdog", "Offline WAV tester (test_audio.py) + live mic sessions"],
                  ["Privacy Shield",    "Scenario simulation — conversation pause edge cases"],
                  ["Cognitive Core",    "Manual prompt evaluation + false-positive suppression review"],
                  ["Full Pipeline",     "Multi-patient end-to-end demo runs — all 7 profiles"],
              ],
              top=Inches(1.82), width=Inches(5.84),
              col_widths=[Inches(1.80), Inches(4.04)])

    # Right table
    section_label(sl, "Integration Quality Measures", Inches(1.55))
    rx = ML + Inches(6.09)
    add_table(sl,
              ["Risk", "Engineering Solution"],
              [
                  ["Alarm fatigue from audio",  "15-sec rolling accumulator — Gemini fires once only"],
                  ["Whisper hallucinations",     "Suppressed on audio clips + entropy transcript filter"],
                  ["Vision/audio bleed",         "Mic silenced during vision segments via is_audio_segment_active"],
                  ["Kaggle API latency",          "Tail-wait polling + 25-second timeout + async tasks"],
                  ["Concurrent Gemini calls",     "gemini_audio_fired one-shot guard per patient"],
                  ["DB row flooding",             "5-second per-sound-type broadcast dedupe"],
              ],
              left=rx, top=Inches(1.82), width=Inches(5.84),
              col_widths=[Inches(2.20), Inches(3.64)])


def s7_results(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Current Results",
                "Evaluated on held-out unseen test sets — no overlap with training data")

    # Fall section
    section_label(sl, "Fall Detection — MoViNet-A2", Inches(1.52))
    sw = Inches(2.85); sx = ML
    fall_stats = [("0.9701", "AUC"), ("86.79%", "Recall"), ("95.83%", "Precision"), ("91.09%", "F1 Score")]
    for val, lbl in fall_stats:
        stat_card(sl, val, lbl, sx, Inches(1.78), width=sw, value_color=SUCCESS)
        sx += sw + Inches(0.13)
    add_text(sl, "950 clips  ·  526 normal  ·  424 fall  ·  Optimal threshold: 0.55",
             ML, Inches(2.84), CW, Inches(0.24), size=Pt(9), color=MID)

    rule(sl, Inches(3.18))

    # Seizure section
    section_label(sl, "Seizure Detection — MoViNet-A2", Inches(3.28))
    sx = ML
    sz_stats = [("0.8162", "AUC"), ("81.72%", "Recall"), ("70.37%", "Precision"), ("75.62%", "F1 Score")]
    for val, lbl in sz_stats:
        stat_card(sl, val, lbl, sx, Inches(3.54), width=sw, value_color=INFO)
        sx += sw + Inches(0.13)
    add_text(sl, "186 clips  ·  93 normal  ·  93 seizure  ·  Optimal threshold: 0.30",
             ML, Inches(4.60), CW, Inches(0.24), size=Pt(9), color=MID)

    rule(sl, Inches(4.94))

    # Audio placeholder
    section_label(sl, "Auditory Watchdog — YAMNet + Faster-Whisper", Inches(5.04))
    rect(sl, ML, Inches(5.30), CW, Inches(0.90), fill_color=RGBColor(0xFF, 0xF7, 0xED))
    rect(sl, ML, Inches(5.30), Pt(4), Inches(0.90), fill_color=WARNING)
    add_text(sl, "Audio Module Metrics — Placeholder (to be completed)",
             ML + Pt(12), Inches(5.36), CW - Pt(16), Inches(0.28),
             size=Pt(10), bold=True, color=WARNING)
    add_text(sl, "Quantitative evaluation metrics to be added here.",
             ML + Pt(12), Inches(5.66), CW - Pt(16), Inches(0.32),
             size=Pt(10), color=MID)


def s8_remaining(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Remaining Work",
                "Items outstanding before final thesis submission")

    row_fills = [
        RGBColor(0xFF, 0xEE, 0xEE),
        RGBColor(0xFF, 0xEE, 0xEE),
        RGBColor(0xFF, 0xF7, 0xED),
        RGBColor(0xFF, 0xF7, 0xED),
        None,
    ]
    add_table(sl,
              ["Priority", "Task", "Details"],
              [
                  ["High",   "System refinement",       "Edge case handling, stress testing under concurrent patients, memory profiling"],
                  ["High",   "Deployment",               "Docker containerisation and production hosting (cloud or on-premises server)"],
                  ["Medium", "Audio module metrics",     "Formal evaluation of YAMNet classification on a labelled respiratory dataset"],
                  ["Medium", "Thesis documentation",     "Final write-up — methodology, evaluation, and results chapters"],
                  ["Low",    "UI polish",                "Accessibility improvements and minor responsive layout review"],
              ],
              top=Inches(1.55), row_h=Inches(0.52),
              row_fills=row_fills,
              col_widths=[Inches(1.10), Inches(2.40), Inches(8.43)])

    # Footer note
    rect(sl, ML, Inches(5.55), CW, Inches(0.75), fill_color=LIGHT_BG)
    add_text(sl,
             "Core functionality is complete and fully demonstrable today. "
             "Remaining items are refinement, formal evaluation, and deployment — not architectural changes.",
             ML + Pt(12), Inches(5.62), CW - Pt(16), Inches(0.60),
             size=Pt(11), color=MID, italic=True)


def s9_demo(prs):
    sl = blank_slide(prs); bg(sl, OFF_WHITE)
    accent_bar(sl)
    slide_title(sl, "Live Demo",
                "Vital Guardian Dashboard — running on localhost:8000")

    section_label(sl, "Demo Sequence", Inches(1.52))
    add_table(sl,
              ["Step", "Patient Profile", "Capability Demonstrated"],
              [
                  ["1", "Fall patient",         "Visual Guardian — fall detection + Gemini binary + clinical enrichment"],
                  ["2", "Seizure patient",       "Visual Guardian — seizure detection + dual-tier Gemini verification"],
                  ["3", "Normal activity",       "Visual Guardian — false positive suppression (Gemini: SUPPRESSED)"],
                  ["4", "Whooping Cough",        "Auditory Watchdog — YAMNet cough detection, rolling accumulator, Gemini audio verdict"],
                  ["5", "Asthma Attack",         "Auditory Watchdog — breathing + wheeze cards consolidated, single Gemini call"],
                  ["6", "Live camera + mic",     "Full real-time pipeline — simultaneous vision and audio on live hardware"],
              ],
              top=Inches(1.78), col_widths=[Inches(0.55), Inches(2.20), Inches(9.18)])

    section_label(sl, "What to Watch For", Inches(4.88))
    cards = [
        ("Risk Gauges",     "Fall and seizure probability arcs update every frame. Watch them spike and trigger the alert."),
        ("Cognitive Core",  "Gemini card populates live — severity badge, clinical narrative, and nursing action checklist."),
        ("Audio Alerts",    "Synthesized alarm tones and Web Speech TTS announce each event without any manual interaction."),
    ]
    cw = Inches(3.84); gap = Inches(0.12); x = ML
    for title, body in cards:
        rect(sl, x, Inches(5.12), cw, Inches(1.14), fill_color=LIGHT_BG)
        add_text(sl, title, x + Pt(10), Inches(5.18), cw - Pt(14), Inches(0.28),
                 size=Pt(10), bold=True, color=ACCENT)
        add_text(sl, body, x + Pt(10), Inches(5.48), cw - Pt(14), Inches(0.68),
                 size=Pt(9.5), color=MID)
        x += cw + gap


def s10_qa(prs):
    sl = blank_slide(prs); bg(sl, WHITE)
    rect(sl, 0, 0, Inches(0.08), H, fill_color=ACCENT)

    add_text(sl, "Thank You", ML, Inches(1.30), CW, Inches(0.70),
             size=Pt(42), bold=True, color=DARK)
    add_text(sl, "Vital Guardian — AI-Based Real-Time Patient Monitoring System",
             ML, Inches(2.08), CW, Inches(0.40), size=Pt(16), color=MID)

    rule(sl, Inches(2.62))

    team = [
        ("Mahad Khan",   "Vision + Dashboard"),
        ("Eman Hassan",  "Cognitive Core"),
        ("Hasan Waqar",  "Auditory Watchdog"),
    ]
    tw = Inches(3.80); tx = ML
    for name, role in team:
        add_text(sl, name, tx, Inches(2.90), tw, Inches(0.32),
                 size=Pt(13), bold=True, color=DARK)
        add_text(sl, role, tx, Inches(3.24), tw, Inches(0.28),
                 size=Pt(11), color=MID)
        tx += tw + Inches(0.12)

    add_text(sl, "Supervisor:  Dr. Basharat Hussain",
             ML, Inches(3.78), CW, Inches(0.30), size=Pt(12), color=DARK)
    add_text(sl, "FAST NUCES Islamabad  ·  Department of Computer Science  ·  2026",
             ML, Inches(4.12), CW, Inches(0.28), size=Pt(11), color=MID)

    rule(sl, Inches(4.55))

    add_text(sl, "Questions & Discussion",
             ML, Inches(4.72), CW, Inches(0.42), size=Pt(20), bold=True, color=DARK)
    add_text(sl,
             "We are happy to elaborate on any module, design decision, evaluation metric, or integration approach.",
             ML, Inches(5.22), CW, Inches(0.36), size=Pt(13), color=MID)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    s1_title(prs);        print("  1/10 — Title")
    s2_overview(prs);     print("  2/10 — Overview")
    s3_feedback(prs);     print("  3/10 — Feedback")
    s4_progress(prs);     print("  4/10 — Progress")
    s5_architecture(prs); print("  5/10 — Architecture")
    s6_testing(prs);      print("  6/10 — Testing")
    s7_results(prs);      print("  7/10 — Results")
    s8_remaining(prs);    print("  8/10 — Remaining")
    s9_demo(prs);         print("  9/10 — Demo")
    s10_qa(prs);          print(" 10/10 — Q&A")

    out = "Vital_Guardian_FYP2_Slides.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")
    print("Open in PowerPoint or Google Slides to present.")

if __name__ == "__main__":
    main()
